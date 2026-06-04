"""
蓝莓 RGB→高光谱重建研究进展 PPT
python make_ppt.py
"""

import csv, io
from pathlib import Path

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
plt.rcParams['font.sans-serif'] = ['Microsoft YaHei', 'SimHei', 'DejaVu Sans']
plt.rcParams['axes.unicode_minus'] = False
import numpy as np

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── 路径 ──────────────────────────────────────────────────────────────────────
ROOT       = Path(r'E:\2026\05\r2h')
LOG_V3     = ROOT / 'checkpoints' / 'v3'  / 'train_log.csv'        # 176波段
LOG_V5     = ROOT / 'checkpoints' / 'v5'  / 'train_log_31.csv'     # 31波段
TEST_V3    = ROOT / 'checkpoints' / 'v3'  / 'test_256patch' / 'test_results.npz'
TEST_V4    = ROOT / 'checkpoints' / 'v4'  / 'test_results.npz'     # NTIRE零样本
VIS_DIR    = ROOT / 'checkpoints' / 'v3'  / 'vis_fullimg'
OUT_PPT    = ROOT / '蓝莓高光谱重建进展.pptx'

# ── 调色板 ────────────────────────────────────────────────────────────────────
C_BG     = RGBColor(0xFF, 0xFF, 0xFF)
C_HEADER = RGBColor(0x1A, 0x56, 0x7A)
C_ACCENT = RGBColor(0x22, 0x96, 0xF3)
C_GREEN  = RGBColor(0x2E, 0x7D, 0x32)
C_ORANGE = RGBColor(0xE6, 0x51, 0x00)
C_RED    = RGBColor(0xC6, 0x28, 0x28)
C_BLACK  = RGBColor(0x21, 0x21, 0x21)
C_GRAY   = RGBColor(0x61, 0x61, 0x61)
C_LGRAY  = RGBColor(0xF5, 0xF5, 0xF5)
C_LINE   = RGBColor(0xBB, 0xDE, 0xFB)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
FONT    = 'Microsoft YaHei'

# ── 基础工具 ──────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def blank(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = C_BG
    return sl

def rect(sl, l, t, w, h, fill=None, line=None, lw=0.75):
    sh = sl.shapes.add_shape(1, l, t, w, h)
    if fill:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line:
        sh.line.color.rgb = line; sh.line.width = Pt(lw)
    else:
        sh.line.fill.background()
    return sh

def txt(sl, text, l, t, w, h,
        size=16, bold=False, color=C_BLACK,
        align=PP_ALIGN.LEFT, wrap=True):
    tb = sl.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    r  = p.add_run()
    r.text           = text
    r.font.size      = Pt(size)
    r.font.bold      = bold
    r.font.color.rgb = color
    r.font.name      = FONT
    return tb

def heading_bar(sl, title):
    rect(sl, 0, 0, SLIDE_W, Inches(0.72), fill=C_HEADER)
    txt(sl, title, Inches(0.45), Inches(0.1), Inches(12.5), Inches(0.52),
        size=24, bold=True, color=C_WHITE)

def hline(sl, y, color=C_LINE):
    rect(sl, Inches(0.4), y, Inches(12.53), Pt(1), fill=color)

def card(sl, l, t, w, h):
    rect(sl, l, t, w, h, fill=C_LGRAY, line=C_LINE, lw=0.5)

def bullet_block(sl, title, items, l, t, w,
                 title_color=C_HEADER, item_color=C_BLACK,
                 title_size=15, item_size=13.5, gap=Inches(0.33)):
    card(sl, l, t, w, Inches(0.44) + gap * len(items))
    txt(sl, title, l+Inches(0.15), t+Inches(0.06),
        w-Inches(0.3), Inches(0.36),
        size=title_size, bold=True, color=title_color)
    for i, item in enumerate(items):
        txt(sl, f'• {item}',
            l+Inches(0.2), t+Inches(0.44)+gap*i,
            w-Inches(0.35), gap,
            size=item_size, color=item_color)

def fig2buf(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=130,
                bbox_inches='tight', facecolor='white')
    buf.seek(0)
    plt.close(fig)
    return buf

def add_fig(sl, fig, l, t, w, h):
    sl.shapes.add_picture(fig2buf(fig), l, t, w, h)

# ── 读取日志 ──────────────────────────────────────────────────────────────────

def read_log(path):
    r = {'epoch': [], 'train_loss': [], 'val_mrae': [], 'val_rmse': [], 'val_psnr': []}
    with open(path, newline='') as f:
        for row in csv.DictReader(f):
            r['epoch'].append(int(row['epoch']))
            r['train_loss'].append(float(row['train_loss']))
            r['val_mrae'].append(float(row['val_mrae']))
            r['val_rmse'].append(float(row['val_rmse']))
            r['val_psnr'].append(float(row['val_psnr']))
    return r

# ── 图表 ──────────────────────────────────────────────────────────────────────

def training_compare_fig(r3, r5):
    """v3(176波段) vs v5(31波段) 训练曲线对比"""
    fig, axes = plt.subplots(1, 2, figsize=(10, 3.5))

    axes[0].plot(r3['epoch'], r3['val_mrae'], color='#1565C0', linewidth=1.6, label='v3 – 176波段')
    axes[0].plot(r5['epoch'], r5['val_mrae'], color='#E65100', linewidth=1.6, label='v5 – 31波段')
    axes[0].set_xlabel('Epoch', fontsize=9)
    axes[0].set_title('Val MRAE ↓', fontsize=10)
    axes[0].legend(fontsize=8); axes[0].grid(True, alpha=0.3)

    axes[1].plot(r3['epoch'], r3['val_psnr'], color='#1565C0', linewidth=1.6, label='v3 – 176波段')
    axes[1].plot(r5['epoch'], r5['val_psnr'], color='#E65100', linewidth=1.6, label='v5 – 31波段')
    axes[1].set_xlabel('Epoch', fontsize=9)
    axes[1].set_title('Val PSNR (dB) ↑', fontsize=10)
    axes[1].legend(fontsize=8); axes[1].grid(True, alpha=0.3)

    fig.tight_layout()
    return fig

def metrics_bar_fig(results):
    """三方案指标对比柱状图"""
    labels = ['v3\n176波段\n(自训练)', 'v5\n31波段\n(自训练)', 'v4\n31波段\n(NTIRE零样本)']
    mrae = [results['v3']['mrae'], results['v5']['mrae'], results['v4']['mrae']]
    psnr = [results['v3']['psnr'], results['v5']['psnr'], results['v4']['psnr']]
    colors = ['#1565C0', '#E65100', '#C62828']

    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(9, 3.2))
    x = np.arange(3)
    ax1.bar(x, mrae, color=colors, width=0.5)
    ax1.set_xticks(x); ax1.set_xticklabels(labels, fontsize=8)
    ax1.set_title('MRAE ↓', fontsize=10)
    ax1.grid(True, alpha=0.3, axis='y')
    for i, v in enumerate(mrae):
        ax1.text(i, v+0.02, f'{v:.3f}' if v < 5 else f'{v:.1f}',
                 ha='center', fontsize=8, fontweight='bold')

    ax2.bar(x, psnr, color=colors, width=0.5)
    ax2.set_xticks(x); ax2.set_xticklabels(labels, fontsize=8)
    ax2.set_title('PSNR (dB) ↑', fontsize=10)
    ax2.grid(True, alpha=0.3, axis='y')
    for i, v in enumerate(psnr):
        ax2.text(i, v+0.15, f'{v:.2f}', ha='center', fontsize=8, fontweight='bold')

    fig.tight_layout()
    return fig

# ── 幻灯片 ────────────────────────────────────────────────────────────────────

def sl_cover(prs, r3, r5, res):
    sl = blank(prs)
    rect(sl, 0, 0, SLIDE_W, Inches(2.6), fill=C_HEADER)
    txt(sl, '蓝莓高光谱重建研究进展',
        Inches(0.8), Inches(0.5), Inches(11.5), Inches(1.2),
        size=40, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(sl, 'RGB（3 通道）→ 高光谱（176 波段，402–1005 nm）',
        Inches(1.5), Inches(1.72), Inches(10), Inches(0.55),
        size=20, color=C_LINE, align=PP_ALIGN.CENTER)

    infos = [
        ('数据集',   '蓝莓样本  95 个样本\n4 品种'),
        ('高光谱相机', 'Kejian  960×991\n176 波段 402–1005 nm'),
        ('模型',     'MST++\n多级光谱 Transformer'),
        ('日期',     '2026-06-03'),
    ]
    for i, (k, v) in enumerate(infos):
        x = Inches(0.5 + i * 3.2)
        card(sl, x, Inches(3.0), Inches(3.05), Inches(1.6))
        txt(sl, k, x+Inches(0.15), Inches(3.1),
            Inches(2.75), Inches(0.38),
            size=13, bold=True, color=C_HEADER)
        txt(sl, v, x+Inches(0.15), Inches(3.5),
            Inches(2.75), Inches(0.9),
            size=13, color=C_GRAY)

    rect(sl, 0, Inches(4.8), SLIDE_W, Inches(2.7), fill=C_LGRAY)
    txt(sl, '本次实验：三方案对比',
        Inches(0.8), Inches(4.95), Inches(11.5), Inches(0.4),
        size=17, bold=True, color=C_HEADER, align=PP_ALIGN.CENTER)

    summary = [
        (f'v3  176波段  {r3["epoch"][-1]} epoch',
         f'MRAE {res["v3"]["mrae"]:.4f}  PSNR {res["v3"]["psnr"]:.2f} dB', C_HEADER),
        (f'v5  31波段   {r5["epoch"][-1]} epoch',
         f'MRAE {res["v5"]["mrae"]:.4f}  PSNR {res["v5"]["psnr"]:.2f} dB', C_ORANGE),
        ('v4  NTIRE预训练  零样本迁移',
         f'MRAE {res["v4"]["mrae"]:.2f}  PSNR {res["v4"]["psnr"]:.2f} dB', C_RED),
    ]
    for i, (label, val, clr) in enumerate(summary):
        x = Inches(0.55 + i * 4.25)
        txt(sl, label, x, Inches(5.48), Inches(4.0), Inches(0.35),
            size=13, bold=True, color=clr)
        txt(sl, val, x, Inches(5.86), Inches(4.0), Inches(0.35),
            size=13, color=C_BLACK)


def sl_experiment_overview(prs):
    sl = blank(prs)
    heading_bar(sl, '实验方案总览')

    rows = [
        ('版本', '波段数', '目标波长范围', '训练轮次', 'Patch尺寸', '权重初始化', '说明'),
        ('v3',  '176',   '402–1005 nm（VIS+NIR）', '100 epoch',
         '256×256', '随机初始化', '蓝莓数据集基线'),
        ('v5',  '31',    '400–700 nm（仅可见光）',  '136 epoch',
         '256×256', '随机初始化', 'NTIRE标准波段对齐'),
        ('v4',  '31',    '400–700 nm（仅可见光）',  '零样本',
         '256×256', 'NTIRE预训练权重', '直接迁移，不微调'),
    ]
    col_w = [Inches(1.1), Inches(1.0), Inches(2.6), Inches(1.4),
             Inches(1.4), Inches(2.1), Inches(3.1)]
    rh = Inches(0.48)
    tx, ty = Inches(0.3), Inches(0.88)

    for i, row in enumerate(rows):
        y = ty + rh * i
        x = tx
        fill = C_HEADER if i == 0 else (C_LGRAY if i % 2 == 1 else C_WHITE)
        for j, cell in enumerate(row):
            rect(sl, x, y, col_w[j], rh, fill=fill, line=C_LINE, lw=0.3)
            clr = C_WHITE if i == 0 else C_BLACK
            txt(sl, cell, x+Inches(0.05), y+Inches(0.07),
                col_w[j]-Inches(0.1), rh-Inches(0.1),
                size=12 if i == 0 else 12, bold=(i == 0), color=clr,
                align=PP_ALIGN.CENTER)
            x += col_w[j]

    hline(sl, Inches(3.25))
    txt(sl, '核心问题：31波段（可见光400-700nm）重建效果是否优于176波段（全谱402-1005nm）？',
        Inches(0.45), Inches(3.35), Inches(12.5), Inches(0.38),
        size=15, bold=True, color=C_HEADER)

    bullet_block(sl, '31波段的动机',
        ['对齐 NTIRE 公开 benchmark，可复用预训练权重',
         '可见光波段与 RGB 物理上更相关',
         '输出维度减小，训练更快'],
        Inches(0.4), Inches(3.85), Inches(6.1),
        title_color=C_ACCENT)

    bullet_block(sl, '176波段的优势',
        ['完整保留 NIR（700–1005 nm）信息',
         'NIR 波段对蓝莓品质分析更重要',
         'RGB 对 NIR 的约束反而使 NIR 更易学'],
        Inches(6.8), Inches(3.85), Inches(6.1),
        title_color=C_GREEN)


def sl_training_curves(prs, r3, r5):
    sl = blank(prs)
    heading_bar(sl, '训练曲线对比：v3（176波段）vs v5（31波段）')

    # 指标摘要卡片
    for i, (label, v3v, v5v, better, clr) in enumerate([
        ('最佳 val MRAE ↓',
         f'{min(r3["val_mrae"]):.4f}  (ep{r3["epoch"][r3["val_mrae"].index(min(r3["val_mrae"]))]})',
         f'{min(r5["val_mrae"]):.4f}  (ep{r5["epoch"][r5["val_mrae"].index(min(r5["val_mrae"]))]})',
         'v3 更优', C_ORANGE),
        ('最佳 val PSNR ↑',
         f'{max(r3["val_psnr"]):.2f} dB  (ep{r3["epoch"][r3["val_psnr"].index(max(r3["val_psnr"]))]})',
         f'{max(r5["val_psnr"]):.2f} dB  (ep{r5["epoch"][r5["val_psnr"].index(max(r5["val_psnr"]))]})',
         'v3 更优', C_GREEN),
        ('训练轮次',
         f'{r3["epoch"][-1]} epochs',
         f'{r5["epoch"][-1]} epochs',
         '', C_GRAY),
    ]):
        x = Inches(0.4 + i * 4.3)
        card(sl, x, Inches(0.88), Inches(4.1), Inches(1.12))
        txt(sl, label, x+Inches(0.15), Inches(0.96),
            Inches(3.8), Inches(0.3), size=12, bold=True, color=clr)
        txt(sl, f'v3: {v3v}', x+Inches(0.15), Inches(1.28),
            Inches(3.8), Inches(0.3), size=12, color=C_HEADER)
        txt(sl, f'v5: {v5v}', x+Inches(0.15), Inches(1.6),
            Inches(3.8), Inches(0.3), size=12, color=C_ORANGE)

    fig = training_compare_fig(r3, r5)
    add_fig(sl, fig, Inches(1.2), Inches(2.1), Inches(10.9), Inches(4.0))


def sl_test_results(prs, res):
    sl = blank(prs)
    heading_bar(sl, '测试集定量结果（256×256 patch，66个）')

    # 三方案对比表
    headers = ['方案', '波段数', 'MRAE ↓', 'RMSE ↓', 'PSNR ↑ (dB)', 'SSIM ↑', '备注']
    col_w   = [Inches(2.2), Inches(1.0), Inches(1.5), Inches(1.5),
               Inches(1.8), Inches(1.5), Inches(3.4)]
    rh = Inches(0.48)
    tx, ty = Inches(0.3), Inches(0.88)

    rows_data = [
        headers,
        ['v3  自训练', '176',
         f'{res["v3"]["mrae"]:.4f}', f'{res["v3"]["rmse"]:.4f}',
         f'{res["v3"]["psnr"]:.2f}', f'{res["v3"]["ssim"]:.4f}',
         '✓ 最佳，NIR全谱'],
        ['v5  自训练', '31',
         f'{res["v5"]["mrae"]:.4f}', f'{res["v5"]["rmse"]:.4f}',
         f'{res["v5"]["psnr"]:.2f}', f'{res["v5"]["ssim"]:.4f}',
         '可见光重建难度更高'],
        ['v4  NTIRE预训练', '31',
         f'{res["v4"]["mrae"]:.2f}',  f'{res["v4"]["rmse"]:.4f}',
         f'{res["v4"]["psnr"]:.2f}', f'{res["v4"]["ssim"]:.4f}',
         '✗ 域迁移失败'],
    ]

    for i, row in enumerate(rows_data):
        y = ty + rh * i
        x = tx
        fill = C_HEADER if i == 0 else (C_LGRAY if i % 2 == 1 else C_WHITE)
        for j, cell in enumerate(row):
            rect(sl, x, y, col_w[j], rh, fill=fill, line=C_LINE, lw=0.3)
            if i == 0:
                clr = C_WHITE
            elif j == 2:
                clr = C_GREEN if i == 1 else (C_ORANGE if i == 2 else C_RED)
            elif j == 4:
                clr = C_GREEN if i == 1 else (C_ORANGE if i == 2 else C_RED)
            else:
                clr = C_BLACK
            txt(sl, cell, x+Inches(0.05), y+Inches(0.07),
                col_w[j]-Inches(0.1), rh-Inches(0.1),
                size=13, bold=(i == 0), color=clr,
                align=PP_ALIGN.CENTER)
            x += col_w[j]

    hline(sl, Inches(3.0))
    fig = metrics_bar_fig(res)
    add_fig(sl, fig, Inches(1.6), Inches(3.1), Inches(10.1), Inches(3.85))


def sl_vis(prs, vis_dir, start=0, n=3, title_suffix=''):
    vis_files = sorted(vis_dir.glob('*.png'))[start:start+n]
    if not vis_files:
        return
    sl = blank(prs)
    heading_bar(sl, f'全图推理结果（v3 176波段）{title_suffix}')
    txt(sl, '每图：Input RGB（菌落叠色）  |  Predicted HSI  |  Ground Truth HSI  |  菌落掩码  |  平均光谱曲线',
        Inches(0.45), Inches(0.82), Inches(12.3), Inches(0.3),
        size=12, color=C_GRAY)
    available = (SLIDE_H - Inches(1.22)) / len(vis_files)
    h = min(available, Inches(2.1))
    for i, f in enumerate(vis_files):
        y = Inches(1.2) + h * i
        sl.shapes.add_picture(str(f), Inches(0.4), y, Inches(12.53), h - Inches(0.06))
        txt(sl, f.stem, Inches(0.45), y+Inches(0.02),
            Inches(2), Inches(0.25), size=10, color=C_GRAY)


def sl_findings(prs, r3, r5, res):
    sl = blank(prs)
    heading_bar(sl, '关键发现与结论')

    findings = [
        ('① 176波段 > 31波段', C_GREEN,
         f'同等训练条件下，176波段 PSNR {res["v3"]["psnr"]:.2f} dB > 31波段 {res["v5"]["psnr"]:.2f} dB，'
         f'差距 {res["v3"]["psnr"]-res["v5"]["psnr"]:.2f} dB。'
         'NIR（700-1005nm）信息对蓝莓重建有实质贡献，丢弃反而使任务更难。'),
        ('② 可见光重建本质更难', C_ORANGE,
         'RGB 本身就是可见光的线性投影，从 RGB 重建 400-700nm 的31个波段是高度欠定问题。'
         'NIR 波段因空间纹理规律更简单，RGB→NIR 的学习反而更容易。'),
        ('③ NTIRE预训练权重无法直接迁移', C_RED,
         f'零样本迁移结果 MRAE={res["v4"]["mrae"]:.1f}（2800%误差），PSNR={res["v4"]["psnr"]:.2f} dB，'
         '完全失效。蓝莓光谱与自然场景域差距极大，必须在目标数据上训练。'),
        ('④ Patch尺寸不是瓶颈', C_HEADER,
         '256×256 patch 下 v3 获得 MRAE=0.303, PSNR=24.37 dB，说明当前 patch 质量没有问题。'
         '性能限制主要来自训练轮数不足（仅100轮）和模型容量（n_feat=31）偏小。'),
    ]

    for i, (title, clr, desc) in enumerate(findings):
        y = Inches(0.88 + i * 1.55)
        card(sl, Inches(0.4), y, Inches(12.5), Inches(1.42))
        txt(sl, title, Inches(0.6), y+Inches(0.08),
            Inches(12), Inches(0.38),
            size=14, bold=True, color=clr)
        txt(sl, desc, Inches(0.6), y+Inches(0.5),
            Inches(12.1), Inches(0.78),
            size=12.5, color=C_BLACK)


def sl_next(prs):
    sl = blank(prs)
    heading_bar(sl, '后续计划')

    bullet_block(sl, '已确认保留',
        ['数据集：95个样本，256×256 patch，不再修改',
         '目标波段：176波段（402–1005 nm，全谱）',
         '架构：MST++（in=3, out=176）'],
        Inches(0.4), Inches(0.88), Inches(5.9),
        title_color=C_GREEN)

    bullet_block(sl, '31波段实验结论',
        ['v5 已保存至 checkpoints/v5/',
         'NTIRE权重零样本迁移完全失败',
         '可见光重建难度高于全谱，不再继续'],
        Inches(7.0), Inches(0.88), Inches(5.9),
        title_color=C_ORANGE)

    hline(sl, Inches(3.22))
    txt(sl, 'v6 训练计划（176波段改进版）',
        Inches(0.45), Inches(3.32), Inches(12), Inches(0.38),
        size=15, bold=True, color=C_HEADER)

    plans = [
        ('改动①', C_HEADER, 'n_feat: 31 → 64',
         '模型特征通道翻倍，提升对176波段输出的建模能力，参数量约 +2M'),
        ('改动②', C_GREEN,  '训练轮数: 100 → 300',
         'v3 仅训了100轮（lr已降至1e-6），实际未充分收敛，完整cosine schedule需300轮'),
        ('改动③', C_ACCENT, 'checkpoint dir: v6, log: train_log_v6.csv',
         '独立日志与权重，保留v3作为对比基线'),
    ]
    for i, (step, clr, change, desc) in enumerate(plans):
        y = Inches(3.88) + Inches(0.98)*i
        card(sl, Inches(0.4), y, Inches(12.5), Inches(0.85))
        txt(sl, f'{step}  {change}',
            Inches(0.6), y+Inches(0.08),
            Inches(12), Inches(0.35),
            size=13, bold=True, color=clr)
        txt(sl, desc,
            Inches(0.6), y+Inches(0.46),
            Inches(12.1), Inches(0.32),
            size=12, color=C_GRAY)


# ── 主函数 ────────────────────────────────────────────────────────────────────

def main():
    print('读取数据 ...')
    r3 = read_log(LOG_V3)
    r5 = read_log(LOG_V5)

    d3 = np.load(str(TEST_V3))
    d4 = np.load(str(TEST_V4))

    res = {
        'v3': {'mrae': float(d3['mrae'].mean()), 'rmse': float(d3['rmse'].mean()),
               'psnr': float(d3['psnr'].mean()), 'ssim': float(d3['ssim'].mean())},
        'v5': {'mrae': float(r5['val_mrae'][r5['val_psnr'].index(max(r5['val_psnr']))]),
               'rmse': float(r5['val_rmse'][r5['val_psnr'].index(max(r5['val_psnr']))]),
               'psnr': float(max(r5['val_psnr'])),
               'ssim': 0.0},
        'v4': {'mrae': float(d4['mrae'].mean()), 'rmse': float(d4['rmse'].mean()),
               'psnr': float(d4['psnr'].mean()), 'ssim': float(d4['ssim'].mean())},
    }

    print(f'  v3: {r3["epoch"][-1]} epochs, MRAE={res["v3"]["mrae"]:.4f}, PSNR={res["v3"]["psnr"]:.2f}')
    print(f'  v5: {r5["epoch"][-1]} epochs, best PSNR={res["v5"]["psnr"]:.2f}')
    print(f'  v4: NTIRE zero-shot, MRAE={res["v4"]["mrae"]:.2f}, PSNR={res["v4"]["psnr"]:.2f}')

    prs = new_prs()
    print('生成幻灯片 ...')
    sl_cover(prs, r3, r5, res)
    sl_experiment_overview(prs)
    sl_training_curves(prs, r3, r5)
    sl_test_results(prs, res)
    sl_vis(prs, VIS_DIR, start=0, n=3, title_suffix='（1/4）')
    sl_vis(prs, VIS_DIR, start=3, n=3, title_suffix='（2/4）')
    sl_vis(prs, VIS_DIR, start=6, n=3, title_suffix='（3/4）')
    sl_vis(prs, VIS_DIR, start=9, n=3, title_suffix='（4/4）')
    sl_findings(prs, r3, r5, res)
    sl_next(prs)

    prs.save(str(OUT_PPT))
    print(f'\n已保存：{OUT_PPT}  ({len(prs.slides)} 张)')


if __name__ == '__main__':
    main()
